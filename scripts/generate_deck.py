"""Generate the Mnemo Lab pptx deck.

Usage:
    pip install python-pptx
    python scripts/generate_deck.py

Output: docs/Mnemo-Lab-Deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------

PRIMARY = RGBColor(0x0A, 0x2E, 0x5C)
ACCENT = RGBColor(0xE3, 0x6A, 0x14)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x21, 0x29, 0x33)
MUTED_TEXT = RGBColor(0x55, 0x5F, 0x6D)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_FG = RGBColor(0xE8, 0xE8, 0xE8)


def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _add_text(slide, *, left, top, width, height, text, size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _add_bullets(slide, *, left, top, width, height, items, size=18, color=DARK_TEXT, level_indent=Pt(14), font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        bullet = "• " if level == 0 else "– "
        run = p.add_run()
        run.text = bullet + text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def _add_bar(slide, *, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_code_block(slide, *, left, top, width, height, code, size=12):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CODE_BG
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(left + Emu(60000), top + Emu(40000), width - Emu(120000), height - Emu(80000))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Consolas"
        run.font.size = Pt(size)
        run.font.color.rgb = CODE_FG


def _header(slide, title: str, *, subtitle: str | None = None):
    _add_bar(slide, left=Inches(0), top=Inches(0), width=Inches(0.18), height=Inches(7.5), color=PRIMARY)
    _add_text(slide, left=Inches(0.5), top=Inches(0.4), width=Inches(12.5), height=Inches(0.7),
              text=title, size=28, bold=True, color=PRIMARY)
    if subtitle:
        _add_text(slide, left=Inches(0.5), top=Inches(1.05), width=Inches(12.5), height=Inches(0.4),
                  text=subtitle, size=16, color=MUTED_TEXT)
    _add_bar(slide, left=Inches(0.5), top=Inches(1.5), width=Inches(0.7), height=Inches(0.04), color=ACCENT)


def _footer(slide, page: int, total: int):
    _add_text(slide, left=Inches(0.5), top=Inches(7.1), width=Inches(6), height=Inches(0.3),
              text="Mnemo Lab — GHCP Advanced Course", size=10, color=MUTED_TEXT)
    _add_text(slide, left=Inches(12), top=Inches(7.1), width=Inches(1), height=Inches(0.3),
              text=f"{page} / {total}", size=10, color=MUTED_TEXT, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_title(prs):
    s = _blank_slide(prs)
    _add_bar(s, left=Inches(0), top=Inches(0), width=Inches(13.333), height=Inches(7.5), color=PRIMARY)
    _add_bar(s, left=Inches(0), top=Inches(5.6), width=Inches(13.333), height=Inches(0.06), color=ACCENT)
    _add_text(s, left=Inches(0.7), top=Inches(2.6), width=Inches(12), height=Inches(1.4),
              text="Mnemo", size=72, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(s, left=Inches(0.7), top=Inches(3.8), width=Inches(12), height=Inches(0.8),
              text="Memoria organizzativa via MCP per il tuo IDE agent",
              size=26, color=RGBColor(0xCB, 0xD6, 0xE4))
    _add_text(s, left=Inches(0.7), top=Inches(4.5), width=Inches(12), height=Inches(0.5),
              text="Asset riusabile · distribuito come repo a sé",
              size=16, color=RGBColor(0xA8, 0xB8, 0xCC))
    _add_text(s, left=Inches(0.7), top=Inches(5.85), width=Inches(12), height=Inches(0.5),
              text="Lab finale — GHCP Advanced", size=18, color=RGBColor(0xE3, 0xE9, 0xF0))
    _add_text(s, left=Inches(0.7), top=Inches(6.3), width=Inches(12), height=Inches(0.4),
              text="Fabio Sabatini", size=14, color=RGBColor(0xCB, 0xD6, 0xE4))
    return s


def slide_problem(prs):
    s = _blank_slide(prs)
    _header(s, "Il problema", subtitle="L'IDE agent è bravo, ma c'è un muro")
    _add_bullets(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.5),
                 items=[
                     "Non conosce le tue spec private — vivono in Wiki, Confluence, Jira, repo separati",
                     "Non ricorda i bug che il tuo team ha già risolto",
                     "È limitato dal training cutoff: la doc che usa potrebbe essere vecchia di mesi",
                     "Risultato: codice plausibile ma scollegato dal contesto reale del progetto",
                     ("Devi raccontargli ogni volta cose che il progetto sa da anni", 1),
                     ("La conoscenza istituzionale non passa al codice nuovo", 1),
                 ], size=20)
    return s


def slide_idea(prs):
    s = _blank_slide(prs)
    _header(s, "L'idea", subtitle="Dai all'IDE agent una memoria organizzativa")
    _add_text(s, left=Inches(0.7), top=Inches(2.2), width=Inches(12), height=Inches(0.6),
              text="«Diamo all'IDE agent una memoria organizzativa.»",
              size=28, bold=True, color=PRIMARY)
    _add_bullets(s, left=Inches(0.7), top=Inches(3.2), width=Inches(12), height=Inches(3.5),
                 items=[
                     "RAG locale come storage della conoscenza",
                     "MCP come canale standard di accesso",
                     "Agenti CLI come alimentatori autonomi",
                     "L'IDE agent (GHCP) come consumatore — invariato",
                 ], size=22)
    return s


def slide_why_rag(prs):
    s = _blank_slide(prs)
    _header(s, "Perché RAG locale", subtitle="Tradeoff espliciti")
    _add_bullets(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.5),
                 items=[
                     "Privacy: spec proprietarie e bug interni non escono dalla rete aziendale",
                     "Costo: zero per query — nessuna API LLM in retrieval",
                     "Latenza: < 100 ms con embeddings locali (fastembed)",
                     "Offline-capable: dopo l'install, niente Internet richiesto",
                     "Versionabile: il `data/` è un artefatto che puoi committare/sincronizzare",
                 ], size=20)
    return s


def slide_why_mcp(prs):
    s = _blank_slide(prs)
    _header(s, "Perché MCP", subtitle="Il \"USB-C\" degli AI tools")
    _add_bullets(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.5),
                 items=[
                     "Standard aperto promosso da Anthropic — vendor-neutral",
                     "Supporto crescente: GitHub Copilot, Claude Code, Cursor, Continue, Zed",
                     "Disaccoppia data source ↔ AI client: scrivi una volta, usi ovunque",
                     "Tool tipizzati (input/output schema) — l'agent ragiona sulla forma, non sulla stringa",
                     "Trasporto stdio o SSE — locale o remoto",
                 ], size=20)
    return s


def slide_architecture(prs):
    s = _blank_slide(prs)
    _header(s, "Architettura", subtitle="Due processi, due lifecycle, comunicano via il RAG store")
    diagram = (
        "  Wiki / ADO / Jira / Confluence       Coding Agent (GHCP in IDE)\n"
        "           │                                       │\n"
        "           │  pull/delta                           │  MCP tools\n"
        "           ▼                                       ▼\n"
        "  ┌──────────────────────┐               ┌──────────────────────┐\n"
        "  │   Agentic CLI        │  upsert       │   MCP Server         │\n"
        "  │   (scheduled)        │ ───────────▶  │   (mnemo, stdio)     │\n"
        "  │   • plan             │               │                      │\n"
        "  │   • extract          │  ◀─────────── │                      │\n"
        "  │   • classify         │   read        └──────────┬───────────┘\n"
        "  │   • cross-ref        │                          │\n"
        "  └──────────────────────┘             ┌──────────────────────────┐\n"
        "                                       │   RAG store (Chroma /    │\n"
        "                                       │   LanceDB) — 2 colls     │\n"
        "                                       │   specs · bug_memory     │\n"
        "                                       └──────────────────────────┘\n"
    )
    _add_code_block(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.7),
                    code=diagram, size=11)
    return s


def slide_two_axes(prs):
    s = _blank_slide(prs)
    _header(s, "Due assi di conoscenza", subtitle="Stesso pattern, due collection")

    # Left: specs
    _add_bar(s, left=Inches(0.7), top=Inches(2.1), width=Inches(5.8), height=Inches(0.07), color=ACCENT)
    _add_text(s, left=Inches(0.7), top=Inches(2.2), width=Inches(5.8), height=Inches(0.5),
              text="specs — cosa va costruito", size=22, bold=True, color=PRIMARY)
    _add_bullets(s, left=Inches(0.7), top=Inches(2.85), width=Inches(5.8), height=Inches(3.5),
                 items=[
                     "User stories, acceptance criteria",
                     "ADR (Architecture Decision Records)",
                     "Epic, BDD scenarios, contratti API",
                     "Sorgente reale: Wiki ADO/Confluence/SharePoint",
                 ], size=17)

    # Right: bugs
    _add_bar(s, left=Inches(6.9), top=Inches(2.1), width=Inches(5.8), height=Inches(0.07), color=ACCENT)
    _add_text(s, left=Inches(6.9), top=Inches(2.2), width=Inches(5.8), height=Inches(0.5),
              text="bug_memory — cosa è andato storto", size=22, bold=True, color=PRIMARY)
    _add_bullets(s, left=Inches(6.9), top=Inches(2.85), width=Inches(5.8), height=Inches(3.5),
                 items=[
                     "Bug risolti, symptom + root cause + fix",
                     "File toccati, PR di fix, pattern tags",
                     "Post-mortem, incident report",
                     "Sorgente reale: ADO Work Items / Jira / GitHub Issues",
                 ], size=17)
    return s


def slide_stack(prs):
    s = _blank_slide(prs)
    _header(s, "Stack tecnico", subtitle="Production-ready, minimo numero di dipendenze")
    _add_bullets(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.5),
                 items=[
                     "Chroma (default) — embedded, persistente, cosine similarity",
                     "LanceDB (opzionale) — moderna, hybrid search nativa",
                     "fastembed — ONNX, no PyTorch, cold start istantaneo",
                     "langchain-text-splitters — chunking token-aware, struttura-aware",
                     "MCP Python SDK + FastMCP — server stdio con decoratori @mcp.tool",
                     "pydantic-settings — config tipata da .env",
                     "PEP 544 Protocols — interfacce strutturali per gli adapter",
                 ], size=18)
    return s


def slide_agents(prs):
    s = _blank_slide(prs)
    _header(s, "I tre agenti custom", subtitle="Due alimentano, uno consuma")
    _add_bullets(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.7),
                 items=[
                     "Agent #1 — Specs ingestion CLI (schedulato)",
                     ("plan: cosa è cambiato dall'ultimo run?", 1),
                     ("extract: parsing markdown, frontmatter YAML, normalizzazione", 1),
                     ("cross-ref: bug ↔ spec, file di codice mentionati", 1),
                     "Agent #2 — Bug ingestion CLI (schedulato)",
                     ("pull: work item risolti da ADO/Jira", 1),
                     ("extract: symptom, root cause, fix, pattern tags", 1),
                     ("classify: bug indicizzabili vs rumore (typo, doc)", 1),
                     "Agent #3 — IDE consumer (GHCP Agent Mode)",
                     ("query: chain di tool MCP per multi-step reasoning", 1),
                     ("synthesize: codice spec-aligned e lesson-aware", 1),
                 ], size=16)
    return s


def slide_demo_overview(prs):
    s = _blank_slide(prs)
    _header(s, "I 4 demo", subtitle="Due implementazioni, due bug — tutti via Mnemo")
    _add_bullets(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.7),
                 items=[
                     "Demo 1 (US-102) — implementare endpoint GET /tasks con filtri",
                     ("scoprire BUG-502 da solo, applicare pattern Sentinel da ADR-002", 1),
                     "Demo 2 (US-302) — migrazione DB indice composito",
                     ("codice che documenta la propria storia, citando BUG-504", 1),
                     "Demo 3 (Bug kanban revert) — diagnosi pattern noto",
                     ("BUG-503 riconosciuto, fix shape proposto, file individuato", 1),
                     "Demo 4 (Performance regression) — bug memory ⨯ stato codice",
                     ("verifica se la migrazione di BUG-504 esiste nel branch corrente", 1),
                 ], size=17)
    return s


def slide_demo_us102(prs):
    s = _blank_slide(prs)
    _header(s, "Demo 1 — US-102", subtitle="Implementare GET /api/v1/tasks con filtri")
    _add_text(s, left=Inches(0.7), top=Inches(1.9), width=Inches(12), height=Inches(0.5),
              text="Prompt: «Implementa US-102. Usa la nostra memoria di progetto.»",
              size=15, color=MUTED_TEXT)
    _add_bullets(s, left=Inches(0.7), top=Inches(2.5), width=Inches(12), height=Inches(3.5),
                 items=[
                     "GHCP → query_specs(\"US-102\") · acceptance criteria + scenari BDD",
                     "GHCP → get_spec(\"ADR-002\") · pattern Sentinel + response shape",
                     "GHCP → query_bugs(\"filter assignee null\") · BUG-502 (filtro perso)",
                     "GHCP genera: endpoint + service con _UNSET + test 3 casi",
                 ], size=18)
    _add_bar(s, left=Inches(0.7), top=Inches(6.0), width=Inches(0.15), height=Inches(0.6), color=ACCENT)
    _add_text(s, left=Inches(0.95), top=Inches(6.05), width=Inches(11), height=Inches(0.6),
              text="WOW: GHCP scopre da solo BUG-502 e applica il pattern. Senza Mnemo: BUG-602.",
              size=15, bold=True, color=ACCENT)
    return s


def slide_demo_us302(prs):
    s = _blank_slide(prs)
    _header(s, "Demo 2 — US-302", subtitle="Migrazione DB con indice composito parziale")
    _add_text(s, left=Inches(0.7), top=Inches(1.9), width=Inches(12), height=Inches(0.5),
              text="Prompt: «Genera la migrazione Alembic per US-302. Commento che spieghi il perché.»",
              size=15, color=MUTED_TEXT)
    _add_bullets(s, left=Inches(0.7), top=Inches(2.5), width=Inches(12), height=Inches(3.5),
                 items=[
                     "GHCP → query_specs(\"US-302\") · AC indice parziale + CONCURRENTLY",
                     "GHCP → query_bugs(\"slow tasks performance\") · BUG-504 con metriche",
                     "GHCP → get_spec(\"ADR-005\") · convenzione \"commenta il perché\"",
                     "GHCP genera: migrazione con docstring che riferisce BUG-504, SLO, dataset",
                 ], size=18)
    _add_bar(s, left=Inches(0.7), top=Inches(6.0), width=Inches(0.15), height=Inches(0.6), color=ACCENT)
    _add_text(s, left=Inches(0.95), top=Inches(6.05), width=Inches(11), height=Inches(0.6),
              text="WOW: il codice generato documenta la propria storia per il prossimo maintainer.",
              size=15, bold=True, color=ACCENT)
    return s


def slide_demo_bug503(prs):
    s = _blank_slide(prs)
    _header(s, "Demo 3 — Bug nuovo: kanban revert", subtitle="Riconoscere un pattern già visto")
    _add_text(s, left=Inches(0.7), top=Inches(1.9), width=Inches(12), height=Inches(0.5),
              text="Bug report: «Card del kanban tornano indietro dopo qualche secondo, no errori»",
              size=15, color=MUTED_TEXT)
    _add_bullets(s, left=Inches(0.7), top=Inches(2.5), width=Inches(12), height=Inches(3.5),
                 items=[
                     "GHCP → query_bugs(\"kanban card revert drag drop\") · BUG-503 top hit",
                     "GHCP → get_bug(\"BUG-503\") · root cause + fix + files + PR #312",
                     "GHCP → get_spec(\"US-203\") · cross-ref ADR-004 \"no fire-and-forget\"",
                     "GHCP apre useTaskMutation.ts → individua useQuickStatusChange → patch",
                 ], size=18)
    _add_bar(s, left=Inches(0.7), top=Inches(6.0), width=Inches(0.15), height=Inches(0.6), color=ACCENT)
    _add_text(s, left=Inches(0.95), top=Inches(6.05), width=Inches(11), height=Inches(0.6),
              text="WOW: bug riconosciuto via memoria del team, non via Google.",
              size=15, bold=True, color=ACCENT)
    return s


def slide_demo_bug504(prs):
    s = _blank_slide(prs)
    _header(s, "Demo 4 — Performance regression /tasks", subtitle="Bug memory ⨯ stato del codice")
    _add_text(s, left=Inches(0.7), top=Inches(1.9), width=Inches(12), height=Inches(0.5),
              text="Bug report: «P99 GET /api/v1/tasks a 3.5s su tenant da 200k task»",
              size=15, color=MUTED_TEXT)
    _add_bullets(s, left=Inches(0.7), top=Inches(2.5), width=Inches(12), height=Inches(3.5),
                 items=[
                     "GHCP → query_bugs(\"slow tasks endpoint\") · BUG-504",
                     "GHCP legge migrations/versions/ per verificare se l'indice esiste",
                     "Ramo A — indice presente: suggerisce EXPLAIN + pg_stat_user_indexes",
                     "Ramo B — indice assente: genera la migrazione (stesso template di US-302)",
                 ], size=18)
    _add_bar(s, left=Inches(0.7), top=Inches(6.0), width=Inches(0.15), height=Inches(0.6), color=ACCENT)
    _add_text(s, left=Inches(0.95), top=Inches(6.05), width=Inches(11), height=Inches(0.6),
              text="WOW: agentic reasoning vero — memoria storica incrociata col branch attuale.",
              size=15, bold=True, color=ACCENT)
    return s


def slide_patterns(prs):
    s = _blank_slide(prs)
    _header(s, "Pattern dimostrati", subtitle="Cosa portarsi a casa dal codice")
    _add_bullets(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.5),
                 items=[
                     "PEP 544 Protocols — interfacce strutturali per VectorStore, Embedder, RagPipeline",
                     "Factory + lazy imports — extras opzionali (lance, llamaindex) senza overhead",
                     "Multi-collection RAG — N assi di conoscenza, stesso scheletro",
                     "Adapter pattern — l'unica cosa che cambia tra lab e produzione è il loader",
                     "Capability detection — SupportsHybridSearch attiva automaticamente lo hybrid path",
                     "pydantic-settings — config tipata da env, validazione gratuita",
                 ], size=18)
    return s


def slide_generalize(prs):
    s = _blank_slide(prs)
    _header(s, "Lo stesso pattern, N domini", subtitle="Dove altro questo schema scala")
    _add_bullets(s, left=Inches(0.7), top=Inches(2.0), width=Inches(12), height=Inches(4.5),
                 items=[
                     "Compliance & policy — regole interne, GDPR, audit",
                     "Runbook & playbook — procedure operative, OnCall",
                     "Incident knowledge — post-mortem aggregati",
                     "Dependency intelligence — changelog librerie usate, breaking changes",
                     "Customer feedback — ticket ricorrenti, pattern di problemi",
                     "Tutti: nuovo adapter di ingestion, MCP server invariato, IDE invariato",
                 ], size=18)
    return s


def slide_takeaways(prs):
    s = _blank_slide(prs)
    _header(s, "Take-aways", subtitle="3 idee da portare a casa")
    _add_text(s, left=Inches(0.7), top=Inches(2.2), width=Inches(12), height=Inches(0.6),
              text="1. MCP è il canale, non il valore",
              size=24, bold=True, color=PRIMARY)
    _add_text(s, left=Inches(0.9), top=Inches(2.8), width=Inches(12), height=Inches(0.5),
              text="Il valore è la conoscenza che ci metti dentro.",
              size=16, color=MUTED_TEXT)

    _add_text(s, left=Inches(0.7), top=Inches(3.6), width=Inches(12), height=Inches(0.6),
              text="2. L'agente è utile quando ha memoria del contesto privato",
              size=24, bold=True, color=PRIMARY)
    _add_text(s, left=Inches(0.9), top=Inches(4.2), width=Inches(12), height=Inches(0.5),
              text="Spec, ADR, bug, decisioni: la conoscenza che vale è quella che Internet non sa.",
              size=16, color=MUTED_TEXT)

    _add_text(s, left=Inches(0.7), top=Inches(5.0), width=Inches(12), height=Inches(0.6),
              text="3. \"Cambia solo l'adapter\"",
              size=24, bold=True, color=PRIMARY)
    _add_text(s, left=Inches(0.9), top=Inches(5.6), width=Inches(12), height=Inches(0.5),
              text="Mnemo è un repo a sé — un pip install, un .env, serve N progetti diversi.",
              size=16, color=MUTED_TEXT)
    return s


def slide_qna(prs):
    s = _blank_slide(prs)
    _add_bar(s, left=Inches(0), top=Inches(0), width=Inches(13.333), height=Inches(7.5), color=PRIMARY)
    _add_text(s, left=Inches(0.7), top=Inches(2.6), width=Inches(12), height=Inches(1.4),
              text="Domande?", size=72, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    _add_text(s, left=Inches(0.7), top=Inches(4.0), width=Inches(12), height=Inches(0.5),
              text="Mnemo (server riusabile): github.com/fsabatini82/mnemo",
              size=20, color=RGBColor(0xCB, 0xD6, 0xE4))
    _add_text(s, left=Inches(0.7), top=Inches(4.55), width=Inches(12), height=Inches(0.5),
              text="Lab materials: session-6-rag-mcp-lab",
              size=20, color=RGBColor(0xCB, 0xD6, 0xE4))
    _add_text(s, left=Inches(0.7), top=Inches(5.3), width=Inches(12), height=Inches(0.5),
              text="MCP spec: modelcontextprotocol.io",
              size=15, color=RGBColor(0xA8, 0xB8, 0xCC))
    _add_text(s, left=Inches(0.7), top=Inches(5.75), width=Inches(12), height=Inches(0.5),
              text="Reference servers: github.com/modelcontextprotocol/servers",
              size=15, color=RGBColor(0xA8, 0xB8, 0xCC))
    return s


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------


def build_deck(output: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_title,
        slide_problem,
        slide_idea,
        slide_why_rag,
        slide_why_mcp,
        slide_architecture,
        slide_two_axes,
        slide_stack,
        slide_agents,
        slide_demo_overview,
        slide_demo_us102,
        slide_demo_us302,
        slide_demo_bug503,
        slide_demo_bug504,
        slide_patterns,
        slide_generalize,
        slide_takeaways,
        slide_qna,
    ]

    total = len(builders)
    for idx, builder in enumerate(builders, start=1):
        slide = builder(prs)
        if idx > 1 and idx < total:
            _footer(slide, idx, total)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


if __name__ == "__main__":
    import sys
    repo_root = Path(__file__).resolve().parents[1]
    out = repo_root / "docs" / "Mnemo-Lab-Deck.pptx"
    try:
        final = build_deck(out)
        print(f"Deck written to: {final}")
    except PermissionError:
        # File is locked (likely open in PowerPoint). Fall back to a
        # versioned filename so the run still produces an artifact.
        fallback = out.with_name(out.stem + ".new.pptx")
        final = build_deck(fallback)
        print(f"Deck locked — wrote fallback to: {final}", file=sys.stderr)
        print("Close the original in PowerPoint and replace it manually.", file=sys.stderr)
